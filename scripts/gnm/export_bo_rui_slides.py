"""
Generate GNM_VLNVerse_Bo_Rui_Meeting.pptx from BO_RUI_SLIDES.md.

Usage:
    python3 scripts/gnm/export_bo_rui_slides.py

Output:
    results/bo_reviewer_packet/exported/GNM_VLNVerse_Bo_Rui_Meeting.pptx

Requires:
    pip install python-pptx
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SOURCE = REPO / "results/bo_reviewer_packet/BO_RUI_SLIDES.md"
OUT_DIR = REPO / "results/bo_reviewer_packet/exported"
OUT_FILE = OUT_DIR / "GNM_VLNVerse_Bo_Rui_Meeting.pptx"

# ── Colours ───────────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GREY = RGBColor(0x44, 0x44, 0x44)
MID_GREY = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)       # deep blue for titles
STATUS_DONE = RGBColor(0x1A, 0x56, 0x2A)   # dark green
STATUS_PARTIAL = RGBColor(0x7D, 0x4B, 0x00) # amber
STATUS_PLANNED = RGBColor(0x4A, 0x4A, 0x8A) # muted purple
CODE_BG = RGBColor(0xF4, 0xF4, 0xF4)

# ── Slide dimensions (widescreen 13.33" x 7.5") ───────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.55)
MARGIN_T = Inches(0.35)
CONTENT_W = Inches(12.23)

TITLE_H = Inches(0.85)
CONTENT_Y = Inches(1.25)
CONTENT_H = Inches(5.95)


# ── Parser ────────────────────────────────────────────────────────────────────

class SlideData:
    def __init__(self, number: int, title: str, raw_body: str, notes: str):
        self.number = number
        self.title = title
        self.raw_body = raw_body.strip()
        self.notes = notes.strip()


def parse_slides(md_text: str) -> list[SlideData]:
    """Extract slide data from the markdown file."""
    slides = []
    # Split on ## Slide N — heading
    pattern = re.compile(r"^## Slide (\d+) — (.+)$", re.MULTILINE)
    matches = list(pattern.finditer(md_text))

    for i, m in enumerate(matches):
        number = int(m.group(1))
        title = m.group(2).strip()

        # Content ends at next ## Slide or end of string
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(md_text)
        chunk = md_text[start:end]

        # Split off speaker notes
        notes_marker = re.search(r"\*\*Speaker notes:\*\*", chunk)
        if notes_marker:
            body = chunk[:notes_marker.start()]
            notes_raw = chunk[notes_marker.end():]
            # Strip surrounding --- separators from notes
            notes = re.sub(r"^[\s\-]*", "", notes_raw).strip()
            notes = re.sub(r"[\s\-]*$", "", notes).strip()
        else:
            body = chunk
            notes = ""

        # Strip leading/trailing --- separators from body
        body = re.sub(r"^[\s\-]*\n", "", body)
        body = re.sub(r"\n[\s\-]*$", "", body)

        slides.append(SlideData(number, title, body.strip(), notes))

    return slides


# ── Block parser ──────────────────────────────────────────────────────────────

class TextBlock:
    def __init__(self, text: str):
        self.text = text


class CodeBlock:
    def __init__(self, lang: str, code: str):
        self.lang = lang
        self.code = code


class TableBlock:
    def __init__(self, rows: list[list[str]], has_header: bool):
        self.rows = rows
        self.has_header = has_header


def parse_body(raw: str) -> list:
    """Parse body text into a list of TextBlock / CodeBlock / TableBlock."""
    blocks = []
    lines = raw.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        # Code fence
        if line.strip().startswith("```"):
            lang = line.strip()[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            blocks.append(CodeBlock(lang, "\n".join(code_lines)))
            i += 1
            continue
        # Table: collect consecutive | lines
        if line.strip().startswith("|"):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            rows = []
            has_header = False
            for tl in table_lines:
                # Skip separator rows (|---|---|)
                if re.match(r"^\s*\|[\s\-\|:]+\|\s*$", tl):
                    has_header = bool(rows)
                    continue
                cells = [c.strip() for c in tl.strip().strip("|").split("|")]
                rows.append(cells)
            if rows:
                blocks.append(TableBlock(rows, has_header))
            continue
        # Normal text (accumulate until we hit a code fence or table or blank line)
        text_lines = []
        while i < len(lines):
            l = lines[i]
            if l.strip().startswith("```") or l.strip().startswith("|"):
                break
            text_lines.append(l)
            i += 1
        chunk = "\n".join(text_lines).strip()
        if chunk:
            blocks.append(TextBlock(chunk))
    return blocks


# ── python-pptx helpers ───────────────────────────────────────────────────────

def set_slide_background(slide, colour: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_title_box(slide, text: str, slide_number: int):
    txBox = slide.shapes.add_textbox(MARGIN_L, MARGIN_T, CONTENT_W, TITLE_H)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    # Strip "Slide N — " prefix that's already in the h2, just use the title
    run.text = text
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = ACCENT
    run.font.name = "Calibri"

    # Thin rule under title: a line shape
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    ln = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        MARGIN_L, MARGIN_T + TITLE_H - Inches(0.05),
        MARGIN_L + CONTENT_W, MARGIN_T + TITLE_H - Inches(0.05),
    )
    ln.line.color.rgb = ACCENT
    ln.line.width = Pt(0.75)


def add_slide_number(slide, n: int, total: int):
    nb = slide.shapes.add_textbox(
        Inches(12.1), Inches(7.1), Inches(1.0), Inches(0.3)
    )
    tf = nb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n} / {total}"
    run.font.size = Pt(9)
    run.font.color.rgb = MID_GREY
    run.font.name = "Calibri"


def _status_colour(text: str) -> RGBColor:
    t = text.upper()
    if "DONE" in t:
        return STATUS_DONE
    if "PARTIAL" in t:
        return STATUS_PARTIAL
    if "PLANNED" in t or "PENDING" in t or "CONFIGURED" in t:
        return STATUS_PLANNED
    return NEAR_BLACK


def _strip_md(text: str) -> str:
    """Strip inline markdown: **bold**, `code`, leading #, leading bullets."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text)
    text = re.sub(r"^\s*[-*]\s+", "• ", text)
    text = re.sub(r"^\s*\d+\.\s+", lambda m: m.group(0).strip() + " ", text)
    return text


def add_text_run(para, text: str, bold=False, italic=False,
                 size_pt=14, colour: RGBColor = None, font="Calibri"):
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size_pt)
    run.font.color.rgb = colour or NEAR_BLACK
    run.font.name = font
    return run


def render_inline(para, raw_line: str, base_size=14):
    """Render a line with inline **bold** and `code` spans."""
    # Split on **...** and `...`
    pattern = re.compile(r"(\*\*[^*]+\*\*|`[^`]+`)")
    parts = pattern.split(raw_line)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            inner = part[2:-2]
            c = _status_colour(inner) if any(
                kw in inner.upper()
                for kw in ("DONE", "PARTIAL", "PLANNED", "PENDING", "CONFIGURED")
            ) else NEAR_BLACK
            add_text_run(para, inner, bold=True, size_pt=base_size, colour=c)
        elif part.startswith("`") and part.endswith("`"):
            add_text_run(para, part[1:-1], size_pt=base_size - 1,
                         colour=DARK_GREY, font="Courier New")
        else:
            add_text_run(para, part, size_pt=base_size)


def add_content_textbox(slide, blocks: list, y=None, h=None):
    """Add the main content text box with mixed formatting."""
    y = y if y is not None else CONTENT_Y
    h = h if h is not None else CONTENT_H

    # We'll lay out text blocks and code/table blocks separately
    # For simplicity, render text and inline code in one text frame,
    # and put code blocks / tables in separate text boxes below.
    # Since we can't perfectly interleave in one frame, we do a single
    # pass rendering everything in one text frame with font changes.

    txBox = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first_para = True

    for block in blocks:
        if isinstance(block, TextBlock):
            lines = block.text.split("\n")
            for raw_line in lines:
                line = raw_line.rstrip()
                if not line:
                    # Empty line: blank paragraph
                    p = tf.add_paragraph() if not first_para else tf.paragraphs[0]
                    if first_para:
                        first_para = False
                    p.space_after = Pt(2)
                    continue

                p = tf.add_paragraph() if not first_para else tf.paragraphs[0]
                if first_para:
                    first_para = False
                p.space_after = Pt(3)

                # Detect heading lines (lines starting with **)
                stripped = line.strip()
                if stripped.startswith("**") and stripped.endswith("**") and \
                        stripped.count("**") == 2:
                    inner = stripped[2:-2]
                    c = _status_colour(inner) if any(
                        kw in inner.upper()
                        for kw in ("DONE", "PARTIAL", "PLANNED", "PENDING",
                                   "CONFIGURED")
                    ) else ACCENT
                    add_text_run(p, inner, bold=True, size_pt=15, colour=c)
                elif stripped.startswith("- ") or stripped.startswith("• "):
                    bullet_text = re.sub(r"^[-•]\s+", "", stripped)
                    add_text_run(p, "  • ", size_pt=13, colour=MID_GREY)
                    render_inline(p, bullet_text, base_size=13)
                elif re.match(r"^\d+\.", stripped):
                    add_text_run(p, "  ", size_pt=13)
                    render_inline(p, stripped, base_size=13)
                else:
                    render_inline(p, line, base_size=13)

        elif isinstance(block, CodeBlock):
            code_lines = block.code.split("\n")
            # Add a visual separator
            p = tf.add_paragraph() if not first_para else tf.paragraphs[0]
            if first_para:
                first_para = False
            p.space_before = Pt(4)

            # Render each code line in Courier New
            for cl in code_lines:
                p = tf.add_paragraph()
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                run = p.add_run()
                # Truncate very long lines
                display_line = cl if len(cl) <= 90 else cl[:87] + "..."
                run.text = "  " + display_line
                run.font.name = "Courier New"
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x2B, 0x2B, 0x60)

            # Trailing space
            p = tf.add_paragraph()
            p.space_before = Pt(4)

        elif isinstance(block, TableBlock):
            # Render table as aligned monospace text
            p = tf.add_paragraph() if not first_para else tf.paragraphs[0]
            if first_para:
                first_para = False
            p.space_before = Pt(4)

            for r_idx, row in enumerate(block.rows):
                p = tf.add_paragraph()
                p.space_before = Pt(1)
                p.space_after = Pt(1)
                row_text = "  " + "   ".join(
                    cell[:35].ljust(36) for cell in row
                )
                run = p.add_run()
                run.text = row_text.rstrip()
                run.font.name = "Courier New"
                run.font.size = Pt(10)
                if r_idx == 0 and block.has_header:
                    run.font.bold = True
                    run.font.color.rgb = ACCENT
                else:
                    run.font.color.rgb = DARK_GREY


def add_notes(slide, notes_text: str):
    """Set the speaker notes for a slide."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    for i, line in enumerate(notes_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.name = "Calibri"


# ── Title slide builder ───────────────────────────────────────────────────────

def build_title_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_background(slide, WHITE)

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "GNM-VLNVerse Baseline"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = "Calibri"

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.7))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Reproducible Isaac Sim pipeline for visual-goal navigation"
    run2.font.size = Pt(20)
    run2.font.color.rgb = DARK_GREY
    run2.font.name = "Calibri"

    # Author
    tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "F. Van Laarhoven — Newcastle University"
    run3.font.size = Pt(14)
    run3.font.color.rgb = MID_GREY
    run3.font.name = "Calibri"

    # Opening line at bottom
    tb4 = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.9))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = (
        "I will answer each point by showing the implementation path: "
        "what is done, how it works, which command proves it, and what is still planned."
    )
    run4.font.size = Pt(13)
    run4.font.italic = True
    run4.font.color.rgb = DARK_GREY
    run4.font.name = "Calibri"

    # Thin bottom rule
    ln = slide.shapes.add_connector(
        1, MARGIN_L, Inches(5.6), MARGIN_L + CONTENT_W, Inches(5.6)
    )
    ln.line.color.rgb = ACCENT
    ln.line.width = Pt(0.5)

    add_notes(slide, (
        "Good morning. This walkthrough maps directly to the questions you raised "
        "after the last meeting. I will go through each one in order, showing what "
        "is implemented, how it works at the code level, and which command you can "
        "run to verify it yourself. Where something is still in progress I will say "
        "so explicitly."
    ))


# ── Main build ────────────────────────────────────────────────────────────────

def build_pptx(slides: list[SlideData]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 is the title slide — build it specially
    build_title_slide(prs)
    total = len(slides) + 1  # +1 for title slide

    for slide_data in slides:
        # Skip slide 1 data — it's covered by the title slide above
        if slide_data.number == 1:
            continue

        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        set_slide_background(slide, WHITE)

        add_title_box(slide, slide_data.title, slide_data.number)
        blocks = parse_body(slide_data.raw_body)
        add_content_textbox(slide, blocks)
        add_slide_number(slide, slide_data.number, total - 1)
        add_notes(slide, slide_data.notes)

    return prs


def main():
    if not SOURCE.exists():
        print(f"ERROR: source not found: {SOURCE}", file=sys.stderr)
        sys.exit(1)

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    md_text = SOURCE.read_text(encoding="utf-8")
    slides = parse_slides(md_text)
    print(f"Parsed {len(slides)} slides from {SOURCE.name}")

    prs = build_pptx(slides)
    prs.save(str(OUT_FILE))
    print(f"Saved: {OUT_FILE}")

    # Append speaker notes to a plain-text backup file
    notes_file = OUT_DIR / "speaker_notes.txt"
    with notes_file.open("w", encoding="utf-8") as f:
        f.write("GNM-VLNVerse Baseline — Speaker Notes\n")
        f.write("=" * 60 + "\n\n")
        for s in slides:
            f.write(f"Slide {s.number}: {s.title}\n")
            f.write("-" * 40 + "\n")
            f.write(s.notes + "\n\n")
    print(f"Notes backup: {notes_file}")


if __name__ == "__main__":
    main()
